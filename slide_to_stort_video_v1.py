import os
import shutil
import asyncio
import subprocess
from pptx import Presentation
from moviepy.editor import *
import win32com.client
from PIL import Image, ImageFilter
import pysrt



# Shorts video format
VIDEO_WIDTH = 1080
VIDEO_HEIGHT = 1920
MAX_SHORTS_DURATION = 60  # 60s limit for Shorts

# =====================================================
# Generate Audio + Subtitles via edge-tts subprocess
# =====================================================
async def generate_audio_with_subs(text, audio_path, subs_path, voice="en-GB-SoniaNeural"):
    cmd = [
        "edge-tts",
        "--voice", voice,
        "--text", text,
        "--write-media", audio_path,
        "--write-subtitles", subs_path,
    ]
    process = await asyncio.create_subprocess_exec(*cmd)
    await process.communicate()

# =====================================================
# Export PPTX slides as images
# =====================================================
def export_slides_as_images(pptx_path, output_dir):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(
        os.path.abspath(pptx_path), WithWindow=False
    )
    presentation.SaveAs(os.path.abspath(output_dir), 17)  # 17 = PNG
    presentation.Close()
    powerpoint.Quit()

# =====================================================
# Create vertical karaoke-style slide clip
# =====================================================
def create_karaoke_clip_vertical(image_path, audio_path, subs_path, duration):
    # Background blurred
    bg_img = Image.open(image_path).convert("RGB")
    bg_img = bg_img.resize((VIDEO_WIDTH, VIDEO_HEIGHT))
    bg_img = bg_img.filter(ImageFilter.GaussianBlur(50))
    bg_img_path = image_path.replace(".JPG", "_bg.jpg")
    bg_img.save(bg_img_path)
    bg = ImageClip(bg_img_path).set_duration(duration)

    # Foreground slide
    fg = ImageClip(image_path).resize(width=VIDEO_WIDTH - 200).set_position("center").set_duration(duration)

    # Audio
    audio_clip = AudioFileClip(audio_path)

    # Subtitles
    subs = pysrt.open(subs_path)
    subtitle_clips = []
    for sub in subs:
        txt = sub.text.strip().replace("\n", " ")
        start = sub.start.ordinal / 1000.0
        end = sub.end.ordinal / 1000.0
        subtitle = (
            TextClip(txt, fontsize=55, color="yellow", stroke_color="black", stroke_width=3,
                     method="caption", size=(VIDEO_WIDTH - 100, None), align="center")
            .set_start(start).set_end(end).set_position(("center", VIDEO_HEIGHT - 200))
        )
        subtitle_clips.append(subtitle)

    final_clip = CompositeVideoClip([bg, fg, *subtitle_clips],
                                    size=(VIDEO_WIDTH, VIDEO_HEIGHT)).set_audio(audio_clip).set_duration(duration)

    return final_clip


# =====================================================
# Main PPTX → Shorts video
# =====================================================
def pptx_to_shorts(pptx_path, output_path="shorts_output.mp4"):
    prs = Presentation(pptx_path)
    working_dir = os.getcwd()
    images_dir = os.path.join(working_dir, "slides")
    os.makedirs(images_dir, exist_ok=True)

    export_slides_as_images(pptx_path, images_dir)

    clips = []
    temp_files = []

    for i, slide in enumerate(prs.slides):
        image_path = os.path.join(images_dir, f"Slide{i+1}.JPG")

        # Extract notes (narration text)
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if not notes_text:
            notes_text = f"This is slide number {i+1}."

        audio_path = os.path.join(working_dir, f"slide_{i+1}.mp3")
        subs_path = os.path.join(working_dir, f"slide_{i+1}.srt")

        # Generate audio + subs with edge-tts
        asyncio.run(generate_audio_with_subs(notes_text, audio_path, subs_path))
        temp_files.extend([audio_path, subs_path])

        # Duration from audio
        audio_clip = AudioFileClip(audio_path)
        duration = audio_clip.duration
        audio_clip.close()

        # Karaoke clip (now using original slide image, no forced resize)
        slide_clip = create_karaoke_clip_vertical(
            image_path, audio_path, subs_path, duration
        )
        clips.append(slide_clip)

    # Final video
    final_video = concatenate_videoclips(clips, method="compose")

    # Trim to Shorts limit
    if final_video.duration > MAX_SHORTS_DURATION:
        final_video = final_video.subclip(0, MAX_SHORTS_DURATION)

    final_video.write_videofile(
        output_path, fps=30, codec="libx264", audio_codec="aac"
    )

    # ✅ Post-export check
    with VideoFileClip(output_path) as v:
        print("✅ YouTube Shorts video saved to", output_path)
        print("   Resolution:", v.size)  # should be (1080, 1920)
        print("   Duration:", v.duration, "seconds")

    # Cleanup
    shutil.rmtree(images_dir, ignore_errors=True)
    for f in temp_files:
        try:
            os.remove(f)
        except:
            pass

# =====================================================
# Example usage
# =====================================================
pptx_to_shorts("short1_JavaSyntax.pptx")
