import os
import shutil
import asyncio
import subprocess
from pptx import Presentation
from moviepy.editor import *
import win32com.client
from PIL import Image
import pysrt

VIDEO_WIDTH = 1280
VIDEO_HEIGHT = 720

# ========== Generate Audio + Subtitles via subprocess ==========
async def generate_audio_with_subs(text, audio_path, subs_path, voice="en-GB-SoniaNeural"):
    cmd = [
        "edge-tts",
        "--voice", voice,
        "--text", text,
        "--write-media", audio_path,
        "--write-subtitles", subs_path
    ]
    process = await asyncio.create_subprocess_exec(*cmd)
    await process.communicate()

# ========== Export PPTX Slides to Images ==========
def export_slides_as_images(pptx_path, output_dir):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
    presentation.SaveAs(os.path.abspath(output_dir), 17)  # 17 = ppSaveAsPNG
    presentation.Close()
    powerpoint.Quit()

# ========== Karaoke Subtitle Overlay ==========
def create_karaoke_clip(image_path, audio_path, subs_path, duration):
    img_clip = ImageClip(image_path).set_duration(duration)
    audio_clip = AudioFileClip(audio_path)

    subs = pysrt.open(subs_path)
    subtitle_clips = []

    for sub in subs:
        txt = sub.text.strip().replace("\n", " ")
        start = sub.start.ordinal / 1000.0
        end = sub.end.ordinal / 1000.0

        subtitle = (
            TextClip(
                txt,
                fontsize=40,
                color="yellow",
                stroke_color="black",
                stroke_width=2,
                method="caption",
                size=(VIDEO_WIDTH - 100, None),
                align="center",
            )
            .set_start(start)
            .set_end(end)
            .set_position(("center", VIDEO_HEIGHT - 100))
        )
        subtitle_clips.append(subtitle)

    final_clip = CompositeVideoClip([img_clip, *subtitle_clips]).set_audio(audio_clip)
    return final_clip

# ========== Main PPTX → Video ==========
def pptx_to_video(pptx_path, output_path="output_video.mp4"):
    prs = Presentation(pptx_path)
    working_dir = os.getcwd()
    images_dir = os.path.join(working_dir, "slides")
    os.makedirs(images_dir, exist_ok=True)

    export_slides_as_images(pptx_path, images_dir)

    clips = []
    temp_files = []

    for i, slide in enumerate(prs.slides):
        image_path = os.path.join(images_dir, f"Slide{i+1}.JPG")

        # Extract notes text
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if not notes_text:
            notes_text = f"This is slide number {i+1}."

        audio_path = os.path.join(working_dir, f"slide_{i+1}.mp3")
        subs_path = os.path.join(working_dir, f"slide_{i+1}.srt")
        resized_img_path = os.path.join(working_dir, f"resized_Slide{i+1}.jpg")

        asyncio.run(generate_audio_with_subs(notes_text, audio_path, subs_path))
        temp_files.extend([audio_path, subs_path])

        # Resize image
        img = Image.open(image_path).convert("RGB")
        img = img.resize((VIDEO_WIDTH, VIDEO_HEIGHT), resample=Image.Resampling.LANCZOS)
        img.save(resized_img_path, "JPEG")
        temp_files.append(resized_img_path)

        # Get duration
        audio_clip = AudioFileClip(audio_path)
        duration = audio_clip.duration
        audio_clip.close()

        slide_clip = create_karaoke_clip(resized_img_path, audio_path, subs_path, duration)
        clips.append(slide_clip)

    final_video = concatenate_videoclips(clips, method="compose")
    final_video.write_videofile(output_path, fps=24)

    print(f"✅ Video saved to {output_path}")

    # Cleanup
    shutil.rmtree(images_dir, ignore_errors=True)
    for f in temp_files:
        try:
            os.remove(f)
        except:
            pass

# Example usage
pptx_to_video("mail_compiling_python_with_notes.pptx")
